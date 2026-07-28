"""Build real, editable .pptx decks from a structured spec.

The 'immense graphics' pipeline (docs/BUSINESS-CAPABILITIES.md), first cut:
branded title/section/content layouts, native charts, big-number slides, an
accent color theme, and speaker notes — all written with python-pptx so the
output is a normal PowerPoint file you can keep editing.

Spec shape (all fields optional unless noted):

    {
      "title": "Q3 Review",            # required
      "subtitle": "Board deck",
      "theme": {"accent": "#2E5BFF", "bg": "#0E1116", "fg": "#FFFFFF"},
      "slides": [
        {"type": "section", "title": "Results"},
        {"title": "Highlights", "bullets": ["...", "..."], "notes": "..."},
        {"title": "Revenue", "chart": {
            "kind": "bar",                       # bar | line | pie
            "categories": ["Q1", "Q2", "Q3"],
            "series": [{"name": "Rev", "values": [10, 14, 19]}]}},
        {"title": "ARR", "big_number": "$4.2M", "caption": "+38% YoY"},
      ]
    }
"""

from __future__ import annotations

from typing import Any


def _rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str.lstrip("#").upper())


class _Theme:
    def __init__(self, theme: dict | None) -> None:
        theme = theme or {}
        self.accent = theme.get("accent", "#2E5BFF")
        self.bg = theme.get("bg", "#0E1116")
        self.fg = theme.get("fg", "#FFFFFF")
        self.muted = theme.get("muted", "#9AA4B2")


def build_deck(spec: dict[str, Any], out_path: str) -> str:
    """Render `spec` to a .pptx at `out_path`. Returns the path."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except Exception as exc:  # pragma: no cover - optional dep
        raise RuntimeError("python-pptx not installed (pip install python-pptx)") from exc

    if not spec.get("title"):
        raise ValueError("deck spec requires a 'title'")

    theme = _Theme(spec.get("theme"))
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    EMU_W, EMU_H = prs.slide_width, prs.slide_height

    def fill_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(color)

    def accent_bar(slide, top, height):
        bar = slide.shapes.add_shape(1, 0, top, Inches(0.18), height)  # rectangle
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(theme.accent)
        bar.line.fill.background()
        return bar

    def text(slide, s, left, top, width, height, size, color, bold=False,
             align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = s
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        return box

    # ---- title slide ----
    s = prs.slides.add_slide(blank)
    fill_bg(s, theme.bg)
    accent_bar(s, Inches(2.6), Inches(2.3))
    text(s, spec["title"], Inches(0.7), Inches(2.6), Inches(11.5), Inches(1.6),
         44, theme.fg, bold=True)
    if spec.get("subtitle"):
        text(s, spec["subtitle"], Inches(0.75), Inches(4.1), Inches(11.5),
             Inches(0.8), 20, theme.muted)

    # ---- content slides ----
    for slide_spec in spec.get("slides", []):
        kind = slide_spec.get("type", "content")
        s = prs.slides.add_slide(blank)
        fill_bg(s, theme.bg)

        if kind == "section":
            accent_bar(s, Inches(3.0), Inches(1.5))
            text(s, slide_spec.get("title", ""), Inches(0.7), Inches(3.1),
                 Inches(11.8), Inches(1.3), 36, theme.fg, bold=True)
            _notes(s, slide_spec.get("notes"))
            continue

        # header for every content slide
        accent_bar(s, Inches(0.55), Inches(0.7))
        text(s, slide_spec.get("title", ""), Inches(0.7), Inches(0.45),
             Inches(11.8), Inches(0.9), 28, theme.fg, bold=True)

        if "big_number" in slide_spec:
            text(s, str(slide_spec["big_number"]), Inches(0.7), Inches(2.4),
                 Inches(11.8), Inches(2.0), 96, theme.accent, bold=True,
                 align=PP_ALIGN.CENTER)
            if slide_spec.get("caption"):
                text(s, slide_spec["caption"], Inches(0.7), Inches(4.6),
                     Inches(11.8), Inches(0.8), 22, theme.muted,
                     align=PP_ALIGN.CENTER)

        elif "chart" in slide_spec:
            _add_chart(s, slide_spec["chart"], theme, Inches, Pt)

        else:  # bullets
            box = s.shapes.add_textbox(Inches(0.7), Inches(1.7),
                                       Inches(11.8), Inches(5.0))
            tf = box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(slide_spec.get("bullets", [])):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = f"•  {bullet}"
                run.font.size = Pt(20)
                run.font.color.rgb = _rgb(theme.fg)
                p.space_after = Pt(10)

        _notes(s, slide_spec.get("notes"))

    prs.save(out_path)
    return out_path


def _notes(slide, note: str | None) -> None:
    if note:
        slide.notes_slide.notes_text_frame.text = note


def _add_chart(slide, chart_spec, theme, Inches, Pt) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    kinds = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }
    xl = kinds.get(chart_spec.get("kind", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)

    data = CategoryChartData()
    data.categories = chart_spec.get("categories", [])
    for series in chart_spec.get("series", []):
        data.add_series(series.get("name", "Series"), series.get("values", []))

    gframe = slide.shapes.add_chart(xl, Inches(0.7), Inches(1.7),
                                    Inches(11.8), Inches(5.2), data)
    chart = gframe.chart
    chart.has_legend = len(chart_spec.get("series", [])) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
